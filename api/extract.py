import io
import re
import requests
from pypdf import PdfReader
from pptx import Presentation
from docx import Document
import openpyxl
from flask import Flask, request, jsonify
from youtube_transcript_api import YouTubeTranscriptApi

app = Flask(__name__)

# Set max content length to 25MB
app.config['MAX_CONTENT_LENGTH'] = 25 * 1024 * 1024

def clean_text(text):
    if not text:
        return ""
    if isinstance(text, bytes):
        text = text.decode('utf-8', errors='ignore')
    if not isinstance(text, str):
        text = str(text)
    
    # Remove script and style tags
    text = re.sub(r'<script\b[^>]*>([\s\S]*?)</script>', '', text, flags=re.I)
    text = re.sub(r'<style\b[^>]*>([\s\S]*?)</style>', '', text, flags=re.I)
    # Remove HTML tags
    text = re.sub(r'<[^>]*>', ' ', text)
    
    # Clean weird spacing and extra newlines
    text = re.sub(r'\s+', ' ', text)
    return text.strip()

def extract_youtube_data(url):
    video_id = None
    if 'youtu.be/' in url:
        video_id = url.split('/')[-1].split('?')[0]
    elif 'youtube.com/watch' in url:
        match = re.search(r'v=([^&]+)', url)
        if match:
            video_id = match.group(1)
    
    if not video_id:
        return "Error: Could not extract Video ID from URL."

    metadata_text = ""
    try:
        # Use oEmbed for metadata (Title, Author)
        oembed_url = f"https://www.youtube.com/oembed?url={url}&format=json"
        resp = requests.get(oembed_url, timeout=5).json()
        metadata_text = f"YOUTUBE_METADATA:\nTitle: {resp.get('title')}\nChannel: {resp.get('author_name')}\n"
    except Exception as e:
        metadata_text = "YOUTUBE_METADATA: (Metadata retrieval failed)\n"

    transcript_text = ""
    try:
        # Fetch transcript in Indonesian or English
        transcript_list = YouTubeTranscriptApi.get_transcript(video_id, languages=['id', 'en'])
        transcript_text = " ".join([t['text'] for t in transcript_list])
    except Exception as e:
        transcript_text = f"(No transcript available for this video. Reason: {str(e)})"

    return f"{metadata_text}\nTRANSCRIPT CONTENT:\n{transcript_text}"

def extract_metadata_heuristics(full_text, filename):
    text_str = str(full_text)
    metadata = {
        "title": filename.rsplit('.', 1)[0].replace("_", " ") if '.' in filename else filename,
        "authors": [],
        "year": "",
        "publisher": "",
        "keywords": [],
        "category": "Original Research",
        "type": "Literature"
    }

    year_match = re.search(r'\b(19|20)\d{2}\b', text_str[:5000])
    if year_match:
        metadata["year"] = year_match.group(0)

    publishers = ["Elsevier", "Springer", "IEEE", "MDPI", "Nature", "Science", "Wiley", "Taylor & Francis", "ACM", "Frontiers", "Sage", "MDPI"]
    for pub in publishers:
        if pub.lower() in text_str[:10000].lower():
            metadata["publisher"] = pub
            break

    return metadata

def process_extracted_text(full_text, title):
    cleaned = clean_text(full_text)
    limit_total = 200000
    limited_text = cleaned[:limit_total]

    metadata = extract_metadata_heuristics(limited_text, title)
    
    ai_snippet = limited_text[:7500]
    chunk_size = 20000
    chunks = [limited_text[i:i+chunk_size] for i in range(0, len(limited_text), chunk_size)][:10]

    return {
        **metadata,
        "aiSnippet": ai_snippet,
        "chunks": chunks,
        "fullText": limited_text
    }

@app.route('/api/extract', methods=['POST'])
def extract():
    try:
        # 1. Handle JSON (Direct URL Extraction, e.g., for YouTube)
        if request.is_json:
            data = request.get_json()
            url = data.get('url')
            if url and ('youtube.com' in url or 'youtu.be' in url):
                result_text = extract_youtube_data(url)
                return jsonify({"status": "success", "data": result_text})
        
        # 2. Handle File Uploads (Standard Logic)
        if 'file' not in request.files:
            # Fallback check for form data if not in request.files
            url = request.form.get('url')
            if url and ('youtube.com' in url or 'youtu.be' in url):
                result_text = extract_youtube_data(url)
                return jsonify({"status": "success", "data": result_text})
            return jsonify({"status": "error", "message": "No file or URL part"}), 400
        
        file = request.files['file']
        if file.filename == '':
            return jsonify({"status": "error", "message": "No file selected"}), 400

        filename_lower = file.filename.lower()
        file_bytes = file.read()
        f = io.BytesIO(file_bytes)
        full_text = ""

        if filename_lower.endswith('.pdf'):
            reader = PdfReader(f)
            for page in reader.pages:
                page_text = page.extract_text()
                if page_text: full_text += page_text + "\n"
        elif filename_lower.endswith('.docx'):
            doc = Document(f)
            full_text = "\n".join([para.text for para in doc.paragraphs])
        elif filename_lower.endswith('.pptx'):
            prs = Presentation(f)
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text: full_text += shape.text + "\n"
        elif filename_lower.endswith('.xlsx'):
            wb = openpyxl.load_workbook(f, data_only=True)
            for sheet in wb.worksheets:
                for row in sheet.iter_rows(values_only=True):
                    full_text += " ".join([str(cell) for cell in row if cell is not None]) + "\n"
        elif filename_lower.endswith(('.txt', '.md', '.csv')):
            full_text = file_bytes.decode('utf-8', errors='ignore')
        else:
            return jsonify({"status": "error", "message": f"Format {filename_lower} not supported."}), 400

        if not full_text.strip():
            return jsonify({"status": "error", "message": "Extracted text is empty."}), 422

        result_data = process_extracted_text(full_text, file.filename)
        return jsonify({"status": "success", "data": result_data})
        
    except Exception as e:
        print(f"Extraction error: {str(e)}")
        return jsonify({"status": "error", "message": str(e)}), 500

if __name__ == '__main__':
    app.run(port=5000)